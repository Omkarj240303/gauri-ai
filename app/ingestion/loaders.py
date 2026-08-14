from pathlib import Path

from langchain_core.documents import Document
from langchain_community.document_loaders import PyPDFLoader


def load_pdf(file_path: Path) -> list[Document]:
    """Load a PDF page-by-page."""
    documents = PyPDFLoader(str(file_path)).load()

    for document in documents:
        document.metadata.update(
            {
                "source_file": file_path.name,
                "file_type": "pdf",
            }
        )

    return documents


def load_text(file_path: Path) -> list[Document]:
    """Load a plain-text file as a single document."""
    text = file_path.read_text(encoding="utf-8")

    return [
        Document(
            page_content=text,
            metadata={
                "source_file": file_path.name,
                "file_type": "txt",
            },
        )
    ]


def load_excel(file_path: Path) -> list[Document]:
    """Load an Excel workbook, preserving sheet names."""
    from openpyxl import load_workbook

    workbook = load_workbook(
        filename=file_path,
        read_only=True,
        data_only=True,
    )

    documents = []

    for sheet in workbook.worksheets:
        rows = []

        for row in sheet.iter_rows(values_only=True):
            values = [
                str(value).strip()
                for value in row
                if value is not None
            ]

            if values:
                rows.append(" | ".join(values))

        if rows:
            documents.append(
                Document(
                    page_content="\n".join(rows),
                    metadata={
                        "source_file": file_path.name,
                        "file_type": "xlsx",
                        "sheet_name": sheet.title,
                    },
                )
            )

    workbook.close()

    return documents


def load_powerpoint(file_path: Path) -> list[Document]:
    """Load PowerPoint slides while preserving slide numbers."""
    from pptx import Presentation

    presentation = Presentation(str(file_path))
    documents = []

    for slide_number, slide in enumerate(
        presentation.slides,
        start=1,
    ):
        texts = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())

        if texts:
            documents.append(
                Document(
                    page_content="\n".join(texts),
                    metadata={
                        "source_file": file_path.name,
                        "file_type": "pptx",
                        "slide_number": slide_number,
                    },
                )
            )

    return documents


def load_file(file_path: Path) -> list[Document]:
    """Load one supported file."""
    extension = file_path.suffix.lower()

    if extension == ".pdf":
        return load_pdf(file_path)

    if extension == ".txt":
        return load_text(file_path)

    if extension == ".xlsx":
        return load_excel(file_path)

    if extension == ".pptx":
        return load_powerpoint(file_path)

    return []